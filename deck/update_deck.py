#!/usr/bin/env python3
"""Edit the ReLoop deck IN PLACE (the working, hand-edited .pptx).

Does two things, without touching build_deck.py (which has diverged / is stale):
  1. Replaces slide 10 (old "where items go" menu) with a combined edge-cases table.
  2. Inserts 4 technical slides (AI Grading, Price Engine, Local Routing, Scale/AWS)
     right before the closing slide.

Run:  cd deck && python3 update_deck.py
"""
import copy
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = "ReLoop-Return-Pipeline (1) (2).pptx"
BAK = "ReLoop-Return-Pipeline (1) (2).bak.pptx"

# --- Amazon-native design tokens (identical to build_deck.py) -------------------
INK      = RGBColor(0x13, 0x1A, 0x22)
NAVY     = RGBColor(0x23, 0x2F, 0x3E)
PANEL    = RGBColor(0x2C, 0x3A, 0x4B)
PANEL2   = RGBColor(0x37, 0x47, 0x57)
ORANGE   = RGBColor(0xFF, 0x99, 0x00)
ORANGE_D = RGBColor(0xEC, 0x72, 0x11)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
FOG      = RGBColor(0xD5, 0xDB, 0xDB)
MUTED    = RGBColor(0x9B, 0xA7, 0xB0)
GREEN    = RGBColor(0x2E, 0xC4, 0x8D)
RED      = RGBColor(0xF2, 0x6D, 0x6D)
CODEBG   = RGBColor(0x0C, 0x12, 0x18)
CODEFG   = RGBColor(0xCF, 0xE8, 0xD6)
CYAN     = RGBColor(0x5D, 0xC8, 0xE8)

FONT = "Arial"
MONO = "Consolas"

prs = Presentation(SRC)
EMU_W, EMU_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# --- helpers (ported from build_deck.py, bound to the loaded prs) ----------------
def new_slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(s, x, y, w, h, color, line=None, line_w=1.0, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, x, y, w, h)
    _set_fill(shp, color)
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if radius:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    shp.shadow.inherit = False
    return shp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.05, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, font, italic) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return tb

def Rn(t, size=14, color=FOG, bold=False, font=FONT, italic=False):
    return (t, size, color, bold, font, italic)

def kicker(s, text):
    txt(s, Inches(0.7), Inches(0.42), Inches(11.5), Inches(0.4),
        [[Rn(text.upper(), 12.5, ORANGE, True, FONT)]])

def accent_bar(s, x=0.7, y=0.34, w=0.55):
    rect(s, Inches(x), Inches(y), Inches(w), Inches(0.06), ORANGE)

def title(s, text, y=0.72, size=30, w=12.2):
    txt(s, Inches(0.7), Inches(y), Inches(w), Inches(1.1),
        [[Rn(text, size, WHITE, True, FONT)]], line_spacing=1.0)

def footer(s, n):
    txt(s, Inches(0.7), Inches(7.08), Inches(9), Inches(0.3),
        [[Rn("ReLoop", 9, ORANGE, True), Rn("  ·  grade at the doorstep, decide before the item moves", 9, MUTED, False)]])
    txt(s, Inches(11.9), Inches(7.08), Inches(0.9), Inches(0.3),
        [[Rn(str(n).zfill(2), 9, MUTED, False)]], align=PP_ALIGN.RIGHT)

def bullets(s, x, y, w, h, items, fs=15, gap=9, lead=ORANGE, headcolor=WHITE):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            head, rest = it
            paras.append([Rn("•  ", fs, lead, True), Rn(head, fs, headcolor, True), Rn(rest, fs, FOG, False)])
        else:
            paras.append([Rn("•  ", fs, lead, True), Rn(it, fs, FOG, False)])
    txt(s, x, y, w, h, paras, space_after=gap, line_spacing=1.12)

def code_panel(s, x, y, w, h, lines, title_txt=None, fs=11):
    rect(s, x, y, w, h, CODEBG, line=PANEL2, line_w=1.0, radius=True)
    dy = y + Inches(0.12)
    for i, c in enumerate((RED, ORANGE, GREEN)):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18 + i*0.22), dy, Inches(0.11), Inches(0.11))
        _set_fill(d, c); d.shadow.inherit = False
    if title_txt:
        txt(s, x + Inches(1.0), dy - Inches(0.02), w - Inches(1.2), Inches(0.24),
            [[Rn(title_txt, 9.5, MUTED, False, MONO)]])
    body_y = y + Inches(0.42)
    tb = s.shapes.add_textbox(x + Inches(0.22), body_y, w - Inches(0.44), h - Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    for i, (t, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1); p.space_before = Pt(0); p.line_spacing = 1.02
        r = p.add_run(); r.text = t
        r.font.size = Pt(fs); r.font.name = MONO; r.font.color.rgb = col; r.font.bold = False
    return tb

def simple_table(s, x, y, rows, col_w, header=True, fs=12, row_h=0.5, header_color=ORANGE_D):
    cy = y
    for ri, row in enumerate(rows):
        cx = x
        rh = Inches(row_h)
        bg = header_color if (header and ri == 0) else (PANEL if ri % 2 else NAVY)
        for ci, cell in enumerate(row):
            cwi = col_w[ci]
            cellbox = rect(s, cx, cy, cwi, rh, bg)
            tf = cellbox.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Pt(8); tf.margin_right = Pt(6)
            tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = cell
            r.font.name = FONT; r.font.size = Pt(fs)
            r.font.bold = (header and ri == 0)
            r.font.color.rgb = WHITE if (header and ri == 0) else FOG
            cx += cwi
        cy += rh
    return cy


def clear_slide(s):
    """Remove every shape from a slide (keeps the background)."""
    for shp in list(s.shapes):
        shp._element.getparent().remove(shp._element)


# ================================================================================
# 1) Rewrite slide 10 -> the edge-cases table
# ================================================================================
s10 = prs.slides[9]
clear_slide(s10)
accent_bar(s10)
kicker(s10, "The tricky cases")
title(s10, "“But what about the weird ones?”")
txt(s10, Inches(0.7), Inches(1.62), Inches(12.0), Inches(0.4),
    [[Rn("Every odd case lands on one of two guarantees: play it safe, or fall back to exactly what Amazon does today. So we never do worse.", 12.5, MUTED, False, FONT, True)]])

edge_rows = [
    ["The tricky case", "The worry", "How ReLoop handles it"],
    ["The AI grades it wrong", "A bad call reaches a buyer",
     "Caught by two human checks (driver + station) before any buyer sees it — worst case is one shelf move"],
    ["Too good to bother (sealed, like-new)", "Wasted on a cheap channel",
     "Straight back on the shelf — skips the pipeline, sells before it loses value"],
    ["The AI isn't sure (blurry photos)", "An over-confident wrong guess",
     "The less sure it is, the safer it plays — low confidence falls back to today's warehouse flow"],
    ["Packaging / accessories missing", "Can't resell as new",
     "Repacked at the local station; a missing ₹300 cable becomes a repair worth +₹1,500"],
    ["Nobody nearby wants it", "The item just sits unsold",
     "Demand was priced in up front; an unsold listing auto-cascades to donate / recycle"],
    ["Customer hides damage / fraud", "Refunding a lie",
     "Refund is held until a human verifies; photo-reuse + capture checks flag fakes"],
    ["Every option loses money (cheap item)", "Costs more to process than it's worth",
     "“Keep it, we'll refund you” — no pickup at all; trusted customers only, never with a fraud signal"],
    ["Counterfeit / hazmat / recalled", "Legal & safety risk",
     "Hard safety rules decide first, never overridden by money — routed to seller or certified disposal"],
]
simple_table(s10, Inches(0.7), Inches(2.18), edge_rows,
             [Inches(3.55), Inches(3.15), Inches(5.6)], fs=11, row_h=0.535)
footer(s10, 10)


# ================================================================================
# 2) Four technical slides (built now, appended, then moved before the closer)
# ================================================================================

# --- Tech 1 — AI grading --------------------------------------------------------
s = new_slide(INK)
accent_bar(s); kicker(s, "Under the hood — the eyes")
title(s, "Our own grading model — not an API call")
bullets(s, Inches(0.7), Inches(1.9), Inches(6.15), Inches(4.4), [
    ("We trained our own eyes. ", "A DINOv2 vision model — not a generic API — so grading can run right at the doorstep, with no cloud round-trip per return."),
    ("One photo, four answers at once. ", "The grade, how sure it is, the exact defects and their severity, and an overall damage score."),
    ("Taught on real + fake damage. ", "Clean Amazon catalog photos, computer-made damage with known labels, and real-world defect datasets."),
    ("Honest, not just confident. ", "It's calibrated — “80% sure” really means 80%. We cut its over-confidence by about 75%."),
], fs=14, gap=13)
# right: what it outputs
rect(s, Inches(7.15), Inches(1.9), Inches(5.45), Inches(2.55), PANEL, radius=True)
txt(s, Inches(7.45), Inches(2.12), Inches(4.9), Inches(0.4), [[Rn("What the model hands back", 14, ORANGE, True)]])
out = [
    ("grade:       70% A · 22% B · 8% C", CODEFG),
    ("confidence:  0.91   (calibrated)", CYAN),
    ("defects:     scratched screen · torn box", CODEFG),
    ("damage:      0.12 overall", CODEFG),
]
rect(s, Inches(7.45), Inches(2.62), Inches(4.85), Inches(1.6), CODEBG, radius=True)
txt(s, Inches(7.7), Inches(2.82), Inches(4.5), Inches(1.3),
    [[Rn(t, 12, c, False, MONO)] for t, c in out], line_spacing=1.22, space_after=1)
# right lower: reference check + flywheel
rect(s, Inches(7.15), Inches(4.6), Inches(5.45), Inches(1.85), PANEL, radius=True)
txt(s, Inches(7.45), Inches(4.8), Inches(4.9), Inches(0.4), [[Rn("Two things an API can't do", 14, ORANGE, True)]])
txt(s, Inches(7.45), Inches(5.25), Inches(4.95), Inches(1.1),
    [[Rn("Reference check: ", 12.5, WHITE, True), Rn("matches the returned item to its original catalog photo + reads the serial — catches swaps & fakes.", 12.5, FOG)],
     [Rn("Flywheel: ", 12.5, WHITE, True), Rn("every human check at the hub is a free training label — it gets better just by doing the job.", 12.5, FOG)]],
    line_spacing=1.14, space_after=5)
footer(s, 15)

# --- Tech 2 — Price engine + priority buyers ------------------------------------
s = new_slide(NAVY)
accent_bar(s); kicker(s, "Under the hood — price & buyers")
title(s, "A living price — and buyers found in priority order")
# left: the price
rect(s, Inches(0.7), Inches(1.9), Inches(5.85), Inches(4.55), PANEL, radius=True)
txt(s, Inches(1.0), Inches(2.12), Inches(5.3), Inches(0.4), [[Rn("The price never sits still", 15, ORANGE, True)]])
rect(s, Inches(1.0), Inches(2.62), Inches(5.25), Inches(0.62), CODEBG, radius=True)
txt(s, Inches(1.15), Inches(2.66), Inches(5.0), Inches(0.55),
    [[Rn("price = base × condition × demand", 11.5, CODEFG, False, MONO)],
     [Rn("        × urgency × category", 11.5, CODEFG, False, MONO)]], line_spacing=1.05, space_after=0)
bullets(s, Inches(1.0), Inches(3.45), Inches(5.3), Inches(3.0), [
    ("Real local demand, ", "a regional demand index refreshed hourly from what buyers actually search, view and buy nearby."),
    ("Racing the clock, ", "urgency climbs as the pickup window closes, so the price adapts on its own."),
    ("A pricing agent, ", "XGBoost + a bandit nudges the listed price toward what similar items are truly selling for — a warehouse queue can't."),
], fs=12.5, gap=10)
# right: cascade
rect(s, Inches(6.75), Inches(1.9), Inches(5.85), Inches(4.55), PANEL, radius=True)
rect(s, Inches(6.75), Inches(1.9), Inches(5.85), Inches(0.12), ORANGE)
txt(s, Inches(7.05), Inches(2.12), Inches(5.3), Inches(0.4), [[Rn("We notify buyers in order — not all at once", 15, ORANGE, True)]])
rect(s, Inches(7.05), Inches(2.66), Inches(5.25), Inches(0.5), CODEBG, radius=True)
txt(s, Inches(7.2), Inches(2.74), Inches(5.0), Inches(0.4),
    [[Rn("score = .30 near + .35 intent + .20 fit + .15 recency", 10.5, CYAN, False, MONO)]])
bullets(s, Inches(7.05), Inches(3.35), Inches(5.35), Inches(3.0), [
    ("Intent leads, ", "who already searched, wishlisted, or bought similar things nearby ranks highest."),
    ("Notify #1 first, ", "if they don't bite in the window, it cascades to #2, #3… automatically — a job that runs every 30 min and survives restarts."),
    ("Fair matches only, ", "within 10 km, same city, and inside each buyer's budget and condition floor."),
], fs=12.5, gap=10)
footer(s, 16)

# --- Tech 3 — Local routing / EV engine -----------------------------------------
s = new_slide(INK)
accent_bar(s); kicker(s, "Under the hood — the decision")
title(s, "One deterministic engine, re-run at every checkpoint")
rect(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(1.05), NAVY, line=ORANGE, line_w=1.25, radius=True)
txt(s, Inches(1.0), Inches(2.05), Inches(11.3), Inches(0.7),
    [[Rn("For every path: ", 14.5, ORANGE, True),
      Rn("value (across every possible grade × how fast it sells) − (shipping + handling + cost of being wrong) − carbon", 14.5, WHITE, True),
      Rn("  →  pick the biggest number.", 14.5, FOG)]], line_spacing=1.16)
cards = [
    ("It thinks in probabilities", "Not one guess — a whole distribution. Restock is punished hard for a wrong “A”; donating barely cares."),
    ("It prices time", "Weeks of warehouse waiting is a real cost line, so “sell it locally now” can beat “ship it and wait.”"),
    ("It prices being wrong", "The cost of a mistake on each path sets how confident the AI must be before that path is even allowed."),
]
for i, (h, b) in enumerate(cards):
    cx = Inches(0.7 + i * 4.05)
    rect(s, cx, Inches(3.15), Inches(3.8), Inches(1.95), PANEL, radius=True)
    rect(s, cx, Inches(3.15), Inches(3.8), Inches(0.1), ORANGE)
    txt(s, cx + Inches(0.28), Inches(3.4), Inches(3.3), Inches(0.5), [[Rn(h, 14.5, ORANGE, True)]])
    txt(s, cx + Inches(0.28), Inches(3.95), Inches(3.3), Inches(1.1), [[Rn(b, 12.5, FOG)]], line_spacing=1.16)
bullets(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.5), [
    ("Safety first. ", "A hard rule ladder (counterfeit, hazmat, recall) runs before any money math and is never overridden."),
    ("Glass box. ", "It's plain arithmetic and clear rules — the AI only writes the one-line reason. And it re-runs at pickup and at the hub as the evidence improves."),
], fs=13.5, gap=10)
footer(s, 17)

# --- Tech 4 — Scalability / AWS -------------------------------------------------
s = new_slide(NAVY)
accent_bar(s); kicker(s, "Built to scale on AWS")
title(s, "How this runs at Amazon scale")
aws_rows = [
    ["The piece", "Rides on AWS"],
    ["The grading model", "SageMaker endpoints + Greengrass at the edge for true doorstep inference; Bedrock vision for narration & fallback"],
    ["Photos & Health Cards", "S3 — with Rekognition redacting faces and addresses right at the upload boundary"],
    ["The decision engine", "Stateless Lambda — a pure, replayable function that scales out flat"],
    ["The return's journey", "Step Functions + EventBridge fire the engine at each checkpoint; item state in DynamoDB"],
    ["Pricing & buyer matching", "Hourly batch demand index, live price on read; EventBridge Scheduler + geo DynamoDB drive the cascade, SNS / Pinpoint send the alerts"],
    ["When the AI's unsure", "Amazon A2I human-review queues; CloudWatch + SageMaker Model Monitor watch for drift"],
    ["Getting smarter", "Hub verdicts stream to a Feature Store / S3 data lake → periodic retraining"],
]
simple_table(s, Inches(0.7), Inches(1.8), aws_rows,
             [Inches(3.1), Inches(8.8)], fs=11.5, row_h=0.565)
rect(s, Inches(0.7), Inches(6.42), Inches(11.9), Inches(0.55), NAVY, line=ORANGE, line_w=1.25, radius=True)
txt(s, Inches(1.0), Inches(6.52), Inches(11.3), Inches(0.4),
    [[Rn("Every piece is stateless or event-driven — it rides Amazon's existing trucks, stations and stores. No new hardware, no new buildings.", 12.5, WHITE, True)]],
    align=PP_ALIGN.CENTER)
footer(s, 18)


# ================================================================================
# 3) Reorder: move the original closing slide (index 14) to the very end
# ================================================================================
sld_lst = prs.slides._sldIdLst
ids = list(sld_lst)
closer = ids[14]          # 15th slide = the "ReLoop / landfill is a design flaw" closer
sld_lst.remove(closer)
sld_lst.append(closer)    # now after the four appended tech slides

shutil.copyfile(SRC, BAK)
prs.save(SRC)
print(f"saved {SRC} with {len(prs.slides._sldIdLst)} slides (backup: {BAK})")
