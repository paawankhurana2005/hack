#!/usr/bin/env python3
"""Generate the ReLoop deck — plain-language, intuitive walkthrough of the idea.

Run:  python3 deck/build_deck.py
Out:  deck/ReLoop-Return-Pipeline.pptx

Voice: explain it to a smart friend. No jargon left unexplained. Follow one
returned item through the whole story.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Amazon-native design tokens ------------------------------------------------
INK      = RGBColor(0x13, 0x1A, 0x22)
NAVY     = RGBColor(0x23, 0x2F, 0x3E)
PANEL    = RGBColor(0x2C, 0x3A, 0x4B)
PANEL2   = RGBColor(0x37, 0x47, 0x57)
ORANGE   = RGBColor(0xFF, 0x99, 0x00)
ORANGE_D = RGBColor(0xEC, 0x72, 0x11)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
FOG      = RGBColor(0xD5, 0xDB, 0xDB)
MUTED    = RGBColor(0x9B, 0xA7, 0xB0)
GREEN    = RGBColor(0x2E, 0xC4, 0x8D)
RED      = RGBColor(0xF2, 0x6D, 0x6D)
CODEBG   = RGBColor(0x0C, 0x12, 0x18)
CODEFG   = RGBColor(0xCF, 0xE8, 0xD6)
CYAN     = RGBColor(0x5D, 0xC8, 0xE8)

FONT = "Arial"
MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# --- helpers --------------------------------------------------------------------
def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(s, x, y, w, h, color, line=None, line_w=1.0, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, x, y, w, h)
    _set_fill(shp, color)
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if radius:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    shp.shadow.inherit = False
    return shp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.05, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, font, italic) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return tb

def Rn(t, size=14, color=FOG, bold=False, font=FONT, italic=False):
    return (t, size, color, bold, font, italic)

def kicker(s, text):
    txt(s, Inches(0.7), Inches(0.42), Inches(11.5), Inches(0.4),
        [[Rn(text.upper(), 12.5, ORANGE, True, FONT)]])

def accent_bar(s, x=0.7, y=0.34, w=0.55):
    rect(s, Inches(x), Inches(y), Inches(w), Inches(0.06), ORANGE)

def title(s, text, y=0.72, size=30, w=12.2):
    txt(s, Inches(0.7), Inches(y), Inches(w), Inches(1.1),
        [[Rn(text, size, WHITE, True, FONT)]], line_spacing=1.0)

def footer(s, n):
    txt(s, Inches(0.7), Inches(7.08), Inches(9), Inches(0.3),
        [[Rn("ReLoop", 9, ORANGE, True), Rn("  ·  grade at the doorstep, decide before the item moves", 9, MUTED, False)]])
    txt(s, Inches(11.9), Inches(7.08), Inches(0.9), Inches(0.3),
        [[Rn(str(n).zfill(2), 9, MUTED, False)]], align=PP_ALIGN.RIGHT)

def bullets(s, x, y, w, h, items, fs=15, gap=9, lead=ORANGE, headcolor=WHITE):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            head, rest = it
            paras.append([Rn("•  ", fs, lead, True), Rn(head, fs, headcolor, True), Rn(rest, fs, FOG, False)])
        else:
            paras.append([Rn("•  ", fs, lead, True), Rn(it, fs, FOG, False)])
    txt(s, x, y, w, h, paras, space_after=gap, line_spacing=1.12)

def code_panel(s, x, y, w, h, lines, title_txt=None, fs=11):
    rect(s, x, y, w, h, CODEBG, line=PANEL2, line_w=1.0, radius=True)
    dy = y + Inches(0.12)
    for i, c in enumerate((RED, ORANGE, GREEN)):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18 + i*0.22), dy, Inches(0.11), Inches(0.11))
        _set_fill(d, c); d.shadow.inherit = False
    if title_txt:
        txt(s, x + Inches(1.0), dy - Inches(0.02), w - Inches(1.2), Inches(0.24),
            [[Rn(title_txt, 9.5, MUTED, False, MONO)]])
    body_y = y + Inches(0.42)
    tb = s.shapes.add_textbox(x + Inches(0.22), body_y, w - Inches(0.44), h - Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    for i, (t, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1); p.space_before = Pt(0); p.line_spacing = 1.02
        r = p.add_run(); r.text = t
        r.font.size = Pt(fs); r.font.name = MONO; r.font.color.rgb = col; r.font.bold = False
    return tb

def screenshot_ph(s, x, y, w, h, cmd, caption):
    rect(s, x, y, w, h, NAVY, line=ORANGE, line_w=1.5, radius=True)
    txt(s, x, y + Inches(0.28), w, Inches(0.4),
        [[Rn("📸  PASTE SCREENSHOT HERE", 13, ORANGE, True, FONT)]], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.3), y + Inches(0.82), w - Inches(0.6), Inches(0.5),
        [[Rn("Run:  ", 11, MUTED, True, FONT), Rn(cmd, 11, CODEFG, False, MONO)]], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.3), y + h - Inches(0.66), w - Inches(0.6), Inches(0.5),
        [[Rn(caption, 10.5, FOG, False, FONT, True)]], align=PP_ALIGN.CENTER)

def simple_table(s, x, y, rows, col_w, header=True, fs=12, row_h=0.5, header_color=ORANGE_D):
    cy = y
    for ri, row in enumerate(rows):
        cx = x
        rh = Inches(row_h)
        bg = header_color if (header and ri == 0) else (PANEL if ri % 2 else NAVY)
        for ci, cell in enumerate(row):
            cwi = col_w[ci]
            cellbox = rect(s, cx, cy, cwi, rh, bg)
            tf = cellbox.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Pt(8); tf.margin_right = Pt(6)
            tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = cell
            r.font.name = FONT; r.font.size = Pt(fs)
            r.font.bold = (header and ri == 0)
            r.font.color.rgb = WHITE if (header and ri == 0) else FOG
            cx += cwi
        cy += rh
    return cy

def step_card(s, x, y, w, h, num, head, body, numcolor=ORANGE):
    rect(s, x, y, w, h, PANEL, radius=True)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.25), Inches(0.6), Inches(0.6))
    _set_fill(circ, numcolor); circ.shadow.inherit = False
    tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
    txt(s, x + Inches(0.25), y + Inches(1.0), w - Inches(0.5), Inches(0.5),
        [[Rn(head, 17, WHITE, True)]])
    txt(s, x + Inches(0.25), y + Inches(1.5), w - Inches(0.5), h - Inches(1.7),
        [[Rn(body, 13, FOG)]], line_spacing=1.16)


# ================================================================================
# 1 — Title
# ================================================================================
s = slide(INK)
rect(s, 0, 0, EMU_W, Inches(0.14), ORANGE)
rect(s, 0, Inches(7.36), EMU_W, Inches(0.14), ORANGE)
txt(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.5),
    [[Rn("AMAZON HACKATHON  ·  THE RETURNS PROBLEM", 14, ORANGE, True)]])
txt(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(1.4),
    [[Rn("ReLoop", 62, WHITE, True)]], line_spacing=1.0)
txt(s, Inches(0.9), Inches(3.4), Inches(11.7), Inches(1.6),
    [[Rn("When you return something, we figure out its ", 27, FOG, False),
      Rn("best next home", 27, ORANGE, True),
      Rn(" right at your front door — ", 27, FOG, False),
      Rn("before", 27, ORANGE, True),
      Rn(" it travels anywhere.", 27, FOG, False)]], line_spacing=1.16)
txt(s, Inches(0.9), Inches(5.7), Inches(11), Inches(0.6),
    [[Rn("The smart layer that sits on top of Amazon's returns — so far less gets wasted.", 17, MUTED, False, FONT, True)]])
txt(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.5),
    [[Rn("Live demo:  ", 12, MUTED, False), Rn("reloop-woad.vercel.app", 12, CYAN, True)]])

# ================================================================================
# 2 — What happens when you return something?
# ================================================================================
s = slide()
accent_bar(s); kicker(s, "Start with something familiar")
title(s, "What actually happens after you click “Return”?")
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.65),
    [[Rn("You hand the package back and forget about it. Behind the scenes, it goes on a surprisingly long trip:", 15.5, FOG)]])
journey = [
    ("Picked up", "a driver collects it from your door"),
    ("Local depot", "it stops at a nearby station — just passing through"),
    ("Sorting center", "trucked to a bigger hub to be batched up"),
    ("Returns center", "hauled hundreds of km away — often 500km+"),
    ("Waits in line", "sits for days or weeks before anyone looks at it"),
    ("Finally graded", "only NOW does a person decide what it's worth"),
]
x0 = 0.7
for i, (h, b) in enumerate(journey):
    cx = Inches(x0 + i * 2.0)
    rect(s, cx, Inches(2.55), Inches(1.82), Inches(1.75), PANEL, radius=True)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.68), Inches(2.72), Inches(0.46), Inches(0.46))
    _set_fill(circ, ORANGE if i == 5 else PANEL2); circ.shadow.inherit = False
    tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1); r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = INK if i == 5 else WHITE; r.font.name = FONT
    txt(s, cx + Inches(0.12), Inches(3.3), Inches(1.6), Inches(0.35),
        [[Rn(h, 12.5, WHITE if i < 5 else ORANGE, True)]], align=PP_ALIGN.CENTER)
    txt(s, cx + Inches(0.12), Inches(3.62), Inches(1.6), Inches(0.65),
        [[Rn(b, 9.5, FOG)]], align=PP_ALIGN.CENTER, line_spacing=1.05)
    if i < 5:
        txt(s, cx + Inches(1.72), Inches(3.15), Inches(0.3), Inches(0.4), [[Rn("›", 22, ORANGE, True)]])
rect(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(1.5), NAVY, line=ORANGE, line_w=1.25, radius=True)
txt(s, Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.1),
    [[Rn("The catch:  ", 18, ORANGE, True),
      Rn("the single most important decision — ", 18, WHITE, True),
      Rn("“what should we do with this item?” — happens ", 18, FOG),
      Rn("dead last", 18, ORANGE, True),
      Rn(", after all the driving is already paid for and weeks of value have quietly leaked away.", 18, FOG)]],
    line_spacing=1.2)
footer(s, 2)

# ================================================================================
# 3 — Why that's a huge, expensive problem
# ================================================================================
s = slide(NAVY)
accent_bar(s); kicker(s, "Why it matters")
title(s, "That long trip is expensive — and most items lose")
txt(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(0.6),
    [[Rn("Returns are enormous, and today's process wrings very little value back out of them.", 15.5, FOG)]])
stats = [
    ("1.2–1.5 billion", "packages returned every year"),
    ("$40–88 billion", "what returns cost the industry, yearly"),
    ("only 10–20%", "of returns ever make it back onto a shelf to be resold"),
    ("the rest", "mostly sold off cheap to liquidators — or thrown away"),
]
for i, (big, lab) in enumerate(stats):
    cx = Inches(0.7 + (i % 2) * 6.05)
    cy = Inches(2.5 + (i // 2) * 1.65)
    rect(s, cx, cy, Inches(5.85), Inches(1.45), PANEL, radius=True)
    rect(s, cx, cy, Inches(0.1), Inches(1.45), ORANGE)
    txt(s, cx + Inches(0.35), cy + Inches(0.2), Inches(5.3), Inches(0.6), [[Rn(big, 27, ORANGE, True)]])
    txt(s, cx + Inches(0.35), cy + Inches(0.82), Inches(5.3), Inches(0.5), [[Rn(lab, 13.5, FOG)]], line_spacing=1.1)
txt(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.7),
    [[Rn("Every wasted item is money burned twice: ", 14.5, ORANGE, True),
      Rn("once on the trucking, and again on the value that could have been recovered if we'd acted sooner.", 14.5, FOG)]],
    line_spacing=1.15)
footer(s, 3)

# ================================================================================
# 4 — The simple idea
# ================================================================================
s = slide()
accent_bar(s); kicker(s, "The idea")
title(s, "So… what if we decided at your front door instead?")
txt(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(1.0),
    [[Rn("The whole problem is that we decide ", 19, FOG),
      Rn("after", 19, ORANGE, True),
      Rn(" the item has travelled. ReLoop flips it: look at the item the moment it's returned, and pick its best next step ", 19, FOG),
      Rn("before it moves an inch.", 19, ORANGE, True)]], line_spacing=1.2)
# before / after contrast
rect(s, Inches(0.7), Inches(3.5), Inches(5.85), Inches(2.9), PANEL, radius=True)
txt(s, Inches(1.0), Inches(3.75), Inches(5.3), Inches(0.4), [[Rn("TODAY", 14, MUTED, True)]])
bullets(s, Inches(1.0), Inches(4.25), Inches(5.3), Inches(2.0), [
    "Ship it far away first",
    "Let it sit and lose value",
    "Then decide what it's worth",
    "Usually too late to do anything clever",
], fs=14, gap=10, lead=RED, headcolor=FOG)
rect(s, Inches(6.85), Inches(3.5), Inches(5.75), Inches(2.9), PANEL, radius=True)
rect(s, Inches(6.85), Inches(3.5), Inches(5.75), Inches(0.12), ORANGE)
txt(s, Inches(7.15), Inches(3.75), Inches(5.3), Inches(0.4), [[Rn("WITH RELOOP", 14, ORANGE, True)]])
bullets(s, Inches(7.15), Inches(4.25), Inches(5.3), Inches(2.0), [
    "Decide at the doorstep, instantly",
    "Keep valuable items local and fresh",
    "Skip the pointless long-distance trip",
    "Only the items that truly need it travel",
], fs=14, gap=10, lead=GREEN, headcolor=WHITE)
footer(s, 4)

# ================================================================================
# 5 — How it works in three steps
# ================================================================================
s = slide(NAVY)
accent_bar(s); kicker(s, "The whole thing in three steps")
title(s, "Snap a photo → the AI decides → the item finds its best home")
step_card(s, Inches(0.7), Inches(1.95), Inches(3.8), Inches(4.4), "1",
          "Look at it",
          "You take a few photos at the door. An AI reads the item's condition — like a really fast, really consistent inspector.")
step_card(s, Inches(4.75), Inches(1.95), Inches(3.8), Inches(4.4), "2",
          "Think it through",
          "A decision engine weighs every option — resell nearby, repair, restock, donate, recycle — and picks the one that recovers the most value for THIS item.")
step_card(s, Inches(8.8), Inches(1.95), Inches(3.8), Inches(4.4), "3",
          "Send it the right way",
          "The item goes straight to its best next home. Two quick human checks along the way make sure the AI got it right.")
footer(s, 5)

# ================================================================================
# 6 — Step 1: the eyes (grading), and honesty about uncertainty
# ================================================================================
s = slide()
accent_bar(s); kicker(s, "Step 1 — the eyes")
title(s, "An AI reads the item's condition from photos")
bullets(s, Inches(0.7), Inches(1.9), Inches(6.1), Inches(3.5), [
    ("Trained to judge condition. ", "It looks at the photos and figures out how good a shape the item is in — brand-new, lightly used, worn, or beyond saving."),
    ("It spots the details. ", "Missing charger? Scratched screen? Torn box? It notices, and lists them."),
    ("And here's the important part — ", "it's honest about how sure it is. Instead of one confident guess, it says things like “probably like-new, maybe lightly used.”"),
], fs=14.5, gap=12)
rect(s, Inches(7.05), Inches(1.9), Inches(5.55), Inches(4.4), PANEL, radius=True)
txt(s, Inches(7.35), Inches(2.15), Inches(5.0), Inches(0.5), [[Rn("Why “how sure” matters", 15, ORANGE, True)]])
txt(s, Inches(7.35), Inches(2.7), Inches(5.0), Inches(1.4),
    [[Rn("A confident wrong guess is dangerous — imagine selling a scratched phone to someone as “like new.” ", 13.5, FOG),
      Rn("By keeping track of its own doubt, the engine can play it safe when the photos aren't clear.", 13.5, WHITE, True)]],
    line_spacing=1.2)
rect(s, Inches(7.35), Inches(4.35), Inches(4.95), Inches(1.75), CODEBG, radius=True)
txt(s, Inches(7.6), Inches(4.55), Inches(4.5), Inches(0.35), [[Rn("what the AI hands over:", 11, MUTED, False, MONO)]])
txt(s, Inches(7.6), Inches(4.95), Inches(4.6), Inches(1.1),
    [[Rn("“70% like-new, 20% lightly used,", 12.5, CODEFG, False, MONO)],
     [Rn(" 10% worn — fairly confident.", 12.5, CODEFG, False, MONO)],
     [Rn(" Note: box looks a bit scuffed.”", 12.5, CYAN, False, MONO)]], line_spacing=1.15, space_after=1)
footer(s, 6)

# ================================================================================
# 7 — Step 2: the brain (plain), glass box
# ================================================================================
s = slide(NAVY)
accent_bar(s); kicker(s, "Step 2 — the brain")
title(s, "It thinks like a smart shopkeeper — and shows its work")
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.9),
    [[Rn("For each item it asks a simple question about every option: ", 15.5, FOG),
      Rn("“how much money do we actually get back this way, after all the costs?”", 15.5, WHITE, True),
      Rn(" — then it picks the winner.", 15.5, FOG)]], line_spacing=1.18)
things = [
    ("What it's worth", "the resale value for the condition it's in"),
    ("The costs to get there", "trucking, handling, repairs — subtracted honestly"),
    ("Who wants it nearby", "is there real demand close by, right now?"),
    ("Time", "things lose value while they sit — that counts as a cost too"),
    ("The planet", "carbon saved by not trucking it far is worth money here"),
    ("Cost of a mistake", "how expensive it'd be to fix a wrong call on this path"),
]
for i, (h, b) in enumerate(things):
    cx = Inches(0.7 + (i % 3) * 4.05)
    cy = Inches(2.75 + (i // 3) * 1.35)
    rect(s, cx, cy, Inches(3.8), Inches(1.2), PANEL, radius=True)
    txt(s, cx + Inches(0.25), cy + Inches(0.16), Inches(3.4), Inches(0.35), [[Rn(h, 13.5, ORANGE, True)]])
    txt(s, cx + Inches(0.25), cy + Inches(0.55), Inches(3.4), Inches(0.6), [[Rn(b, 11.5, FOG)]], line_spacing=1.1)
rect(s, Inches(0.7), Inches(5.65), Inches(11.9), Inches(0.95), NAVY, line=ORANGE, line_w=1.25, radius=True)
txt(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.6),
    [[Rn("No black box. ", 15, ORANGE, True),
      Rn("It's plain arithmetic and clear rules — you can read exactly why it chose what it chose. The AI only writes the one-line explanation; the ", 14.5, FOG),
      Rn("rules do the deciding.", 14.5, WHITE, True)]], line_spacing=1.15)
footer(s, 7)

# ================================================================================
# 8 — See the brain think (real example, plain)
# ================================================================================
s = slide()
accent_bar(s); kicker(s, "See it in action")
title(s, "A real example: a returned wireless speaker")
txt(s, Inches(0.7), Inches(1.65), Inches(6.2), Inches(0.9),
    [[Rn("Sealed box, “changed my mind.” The engine lays out every option with real rupee values and picks the best. Here, ", 13, FOG),
      Rn("putting it straight back on the shelf", 13, WHITE, True),
      Rn(" wins.", 13, FOG)]], line_spacing=1.16)
ev = [
    ("Best option chosen:  put it back on sale", ORANGE),
    ("", FOG),
    ("Back on the shelf     +₹1,956   ← winner", GREEN),
    ("Resell nearby         +₹1,952", CODEFG),
    ("Repair & resell       +₹1,734", CODEFG),
    ("Sell in a graded lot    +₹893", CODEFG),
    ("Donate                  +₹182", CODEFG),
    ("", FOG),
    ("Ship it far away        −₹465   ← today's", RED),
    ("                                  default", RED),
]
code_panel(s, Inches(0.7), Inches(2.7), Inches(6.2), Inches(3.7), ev,
           title_txt="the engine's math, in rupees", fs=12.5)
screenshot_ph(s, Inches(7.15), Inches(1.7), Inches(5.45), Inches(2.75),
              "curl /api/route  (the speaker)",
              "The engine's full option-by-option breakdown, live.")
rect(s, Inches(7.15), Inches(4.7), Inches(5.45), Inches(1.7), PANEL, radius=True)
txt(s, Inches(7.4), Inches(4.9), Inches(5.0), Inches(0.4), [[Rn("The punchline", 14, ORANGE, True)]])
txt(s, Inches(7.4), Inches(5.35), Inches(5.0), Inches(1.0),
    [[Rn("Today's default — shipping it far away — is the ", 13.5, FOG),
      Rn("worst", 13.5, WHITE, True),
      Rn(" option here (it loses money). Every smarter choice is money we get to keep.", 13.5, FOG)]],
    line_spacing=1.18)
footer(s, 8)

# ================================================================================
# 9 — What if the AI is unsure or wrong? (safety net)
# ================================================================================
s = slide(NAVY)
accent_bar(s); kicker(s, "The obvious worry")
title(s, "“But what if the AI gets it wrong?”")
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.6),
    [[Rn("Fair question. Three things make a wrong call cheap and rare — not a disaster.", 15.5, FOG)]])
safety = [
    ("When unsure, it plays safe", "The less confident the AI is, the more cautious the option it's allowed to pick. Shaky photos never lead to a bold, risky choice."),
    ("Two humans double-check", "A 30-second look by the pickup driver, then a 10-minute check at the local station. If a person disagrees, the plan updates on the spot — before any buyer sees the item."),
    ("Worst case = today", "If nothing local makes sense, it just does exactly what Amazon does now: send it up the chain. So we can only do better, never worse."),
]
for i, (h, b) in enumerate(safety):
    cy = Inches(2.5 + i * 1.35)
    rect(s, Inches(0.7), cy, Inches(11.9), Inches(1.2), PANEL, radius=True)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), cy + Inches(0.32), Inches(0.55), Inches(0.55))
    _set_fill(circ, ORANGE); circ.shadow.inherit = False
    tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1); r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
    txt(s, Inches(1.75), cy + Inches(0.16), Inches(10.5), Inches(0.4), [[Rn(h, 15.5, ORANGE, True)]])
    txt(s, Inches(1.75), cy + Inches(0.58), Inches(10.6), Inches(0.55), [[Rn(b, 13, FOG)]], line_spacing=1.12)
footer(s, 9)

# ================================================================================
# 10 — the live re-route moment
# ================================================================================
s = slide()
accent_bar(s); kicker(s, "The best part to watch")
title(s, "When a human disagrees, the plan changes instantly")
txt(s, Inches(0.7), Inches(1.7), Inches(6.1), Inches(1.5),
    [[Rn("At the local station, a worker opens the “sealed” box and finds the seal broken. They tap that in — and the engine ", 14.5, FOG),
      Rn("re-decides on the spot.", 14.5, WHITE, True),
      Rn(" “Put it back on the shelf” is off the table, so it instantly switches to “resell nearby.”", 14.5, FOG)]],
    line_spacing=1.2)
flow = [
    ("BEFORE — box looked sealed", MUTED),
    ("   plan: put it back on the shelf", GREEN),
    ("", FOG),
    ("Worker: “actually, seal's broken”", CYAN),
    ("", FOG),
    ("AFTER — engine re-decides instantly", MUTED),
    ("   new plan: resell it nearby", ORANGE),
    ("", FOG),
    ("The mistake was caught before any", CODEFG),
    ("buyer saw it. Fixing it cost one", CODEFG),
    ("shelf move — not a lost parcel.", CODEFG),
]
code_panel(s, Inches(0.7), Inches(3.35), Inches(6.1), Inches(3.15), flow,
           title_txt="live re-route", fs=12)
screenshot_ph(s, Inches(7.15), Inches(1.85), Inches(5.45), Inches(2.6),
              "Hub screen → flip the seal",
              "Watch the decision switch on screen in real time.")
rect(s, Inches(7.15), Inches(4.7), Inches(5.45), Inches(1.8), PANEL, radius=True)
txt(s, Inches(7.4), Inches(4.9), Inches(5.0), Inches(0.4), [[Rn("Bonus: it gets smarter over time", 14, ORANGE, True)]])
txt(s, Inches(7.4), Inches(5.35), Inches(5.0), Inches(1.1),
    [[Rn("Every time a human confirms or corrects the AI, that becomes a free training example. ", 13.5, FOG),
      Rn("The AI's eyes improve just by doing the job.", 13.5, WHITE, True)]],
    line_spacing=1.18)
footer(s, 10)

# ================================================================================
# 11 — where items end up (menu)
# ================================================================================
s = slide(NAVY)
accent_bar(s); kicker(s, "The happy endings")
title(s, "Where your item can go — instead of the landfill")
rows = [
    ["The good next home", "When it's picked", "Why it's better"],
    ["Back on the shelf", "Sealed / like-new, still wanted", "Sold again quickly — no long trip, no weeks lost"],
    ["Resold to someone nearby", "Good condition, buyers close by", "Recovers far more value; cash back in days"],
    ["Repaired, then resold", "One cheap fix makes it worth more", "A ₹300 cable can add ₹1,500 of value"],
    ["Bundled into a graded lot", "Cheaper or uncertain items", "Sold in clearly-labelled boxes that fetch more"],
    ["Donated", "Nobody's buying, but it's usable", "Goodwill + tax credit instead of the bin"],
    ["Recycled", "Truly finished, or legally required", "Materials recovered, done responsibly"],
    ["Thrown away", "Legally the only option", "The number we're trying to SHRINK — every case logged"],
]
cw = [Inches(3.5), Inches(3.9), Inches(4.5)]
simple_table(s, Inches(0.7), Inches(1.75), rows, cw, fs=11.5, row_h=0.585)
footer(s, 11)

# ================================================================================
# 12 — trust / Health Card
# ================================================================================
s = slide()
accent_bar(s); kicker(s, "The trust question — would a stranger buy it?")
title(s, "Every item carries a “health record” — like CARFAX for a car")
txt(s, Inches(0.7), Inches(1.8), Inches(6.1), Inches(3.5),
    [[Rn("People happily buy used cars from strangers because of CARFAX — a trusted history report. Used stuff online is scary because there's no such thing.", 15, FOG)]],
    line_spacing=1.22)
bullets(s, Inches(0.7), Inches(3.3), Inches(6.1), Inches(3.0), [
    ("A Health Card ", "travels with each item: its verified condition, photos, and every check it passed."),
    ("It's Amazon-backed, ", "not a random seller's word — so a buyer can actually trust it."),
    ("And it follows the item ", "to its next owner, and the one after that."),
], fs=14, gap=11)
rect(s, Inches(7.05), Inches(1.9), Inches(5.55), Inches(4.4), PANEL, radius=True)
rect(s, Inches(7.05), Inches(1.9), Inches(5.55), Inches(0.12), ORANGE)
txt(s, Inches(7.35), Inches(2.2), Inches(5.0), Inches(0.45), [[Rn("PRODUCT HEALTH CARD", 14, ORANGE, True)]])
card_lines = [
    ("✓ Condition: like-new (verified)", GREEN),
    ("✓ Checked at pickup", GREEN),
    ("✓ Checked at the local station", GREEN),
    ("✓ Genuine — serial matches", GREEN),
    ("• Photos included", FOG),
    ("• Owner 2 of 2 · CO₂ saved so far", FOG),
]
yy = Inches(2.85)
for t, c in card_lines:
    txt(s, Inches(7.5), yy, Inches(4.9), Inches(0.4), [[Rn(t, 14, c, False)]])
    yy += Inches(0.55)
footer(s, 12)

# ================================================================================
# 13 — it keeps working on its own
# ================================================================================
s = slide(NAVY)
accent_bar(s); kicker(s, "It doesn't stop at “list it”")
title(s, "A little agent keeps working until the item sells")
bullets(s, Inches(0.7), Inches(1.95), Inches(6.1), Inches(4.2), [
    ("It watches each local listing ", "and adjusts the price toward what similar items are actually selling for — something a warehouse queue can never do."),
    ("It finds the right buyers ", "by tapping into who nearby already searched for, wishlisted, or bought similar things."),
    ("It knows when to give up gracefully. ", "If resale truly isn't working, it doesn't let the item rot — it hands it back to the engine, which donates or recycles it instead."),
], fs=15, gap=13)
rect(s, Inches(7.05), Inches(1.95), Inches(5.55), Inches(4.2), PANEL, radius=True)
txt(s, Inches(7.35), Inches(2.2), Inches(5.0), Inches(0.5), [[Rn("Where a buyer sees it", 15, ORANGE, True)]])
txt(s, Inches(7.35), Inches(2.75), Inches(5.0), Inches(1.5),
    [[Rn("Right on the normal product page: ", 13.5, FOG),
      Rn("“Open-box near you — checked, guaranteed, X% off, delivered today.”", 13.5, WHITE, True),
      Rn("  With a gentle nudge if it's already on their wish list.", 13.5, FOG)]],
    line_spacing=1.2)
rect(s, Inches(7.35), Inches(4.45), Inches(4.95), Inches(1.5), CODEBG, radius=True)
txt(s, Inches(7.6), Inches(4.65), Inches(4.5), Inches(1.2),
    [[Rn("A returns warehouse can't", 13, CODEFG, False, MONO)],
     [Rn("negotiate a price or chase a", 13, CODEFG, False, MONO)],
     [Rn("buyer. This little agent does —", 13, CYAN, False, MONO)],
     [Rn("every single day.", 13, CYAN, False, MONO)]], line_spacing=1.12, space_after=1)
footer(s, 13)

# ================================================================================
# 14 — where the money comes from
# ================================================================================
s = slide()
accent_bar(s); kicker(s, "Follow the money")
title(s, "Why this actually saves and makes money")
money = [
    ("Skip the long trip", "For local items, we delete the sorting hop, the 500km+ haul, and the weeks of waiting. That's pure saved cost."),
    ("Sell before value drops", "Electronics and fashion lose value fast. Selling in days instead of weeks keeps more of the price."),
    ("Repairs that pay off", "A tiny fix can bump an item up a grade — the engine only does it when the math clearly works."),
    ("Better liquidation", "Clearly-labelled, graded boxes sell for more than “mystery pallets,” because buyers can see what's inside."),
]
for i, (h, b) in enumerate(money):
    cx = Inches(0.7 + (i % 2) * 6.05)
    cy = Inches(1.95 + (i // 2) * 2.15)
    rect(s, cx, cy, Inches(5.85), Inches(1.9), PANEL, radius=True)
    rect(s, cx, cy, Inches(0.1), Inches(1.9), GREEN)
    txt(s, cx + Inches(0.35), cy + Inches(0.22), Inches(5.3), Inches(0.45), [[Rn(h, 16, GREEN, True)]])
    txt(s, cx + Inches(0.35), cy + Inches(0.75), Inches(5.3), Inches(1.0), [[Rn(b, 13.5, FOG)]], line_spacing=1.18)
footer(s, 14)

# ================================================================================
# 15 — honest homework + returnless
# ================================================================================
s = slide(NAVY)
accent_bar(s); kicker(s, "We did our homework")
title(s, "Grounded in real numbers — even the “just keep it” case")
txt(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(0.6),
    [[Rn("We didn't invent the economics — every figure comes from how Amazon's returns actually work.", 15, FOG)]])
facts = [
    ("Liquidation really pays cents", "Sellers get back only ~5–10% of an item's price through bulk liquidation. We priced against that reality, not a fantasy."),
    ("Processing can cost more than the item", "It can cost ~$27 to process a $100 return. For cheap items, every option loses money — so we handle that too."),
]
for i, (h, b) in enumerate(facts):
    cy = Inches(2.4 + i * 1.15)
    rect(s, Inches(0.7), cy, Inches(6.6), Inches(1.0), PANEL, radius=True)
    txt(s, Inches(0.95), cy + Inches(0.13), Inches(6.1), Inches(0.4), [[Rn(h, 13.5, ORANGE, True)]])
    txt(s, Inches(0.95), cy + Inches(0.52), Inches(6.15), Inches(0.45), [[Rn(b, 11.5, FOG)]], line_spacing=1.08)
rect(s, Inches(7.55), Inches(2.4), Inches(5.05), Inches(3.5), PANEL, radius=True)
rect(s, Inches(7.55), Inches(2.4), Inches(5.05), Inches(0.12), ORANGE)
txt(s, Inches(7.85), Inches(2.7), Inches(4.5), Inches(0.5), [[Rn("The clever last resort", 15, ORANGE, True)]])
txt(s, Inches(7.85), Inches(3.25), Inches(4.5), Inches(2.4),
    [[Rn("Sometimes the smartest move is to just let you keep the item and refund you — no pickup at all.", 14, WHITE, True)],
     [Rn("", 6, FOG)],
     [Rn("It's a real Amazon tactic. We only do it for low-value items, from trusted customers, with no fraud signs — never as a loophole.", 13, FOG)]],
    line_spacing=1.2, space_after=6)
txt(s, Inches(0.7), Inches(4.9), Inches(6.6), Inches(1.0),
    [[Rn("The point: ", 13.5, ORANGE, True),
      Rn("the engine is honest about what each path is really worth — so it never talks itself into a losing move.", 13.5, FOG)]],
    line_spacing=1.16)
footer(s, 15)

# ================================================================================
# 16 — is it real?
# ================================================================================
s = slide()
accent_bar(s); kicker(s, "Is any of this real?")
title(s, "Yes — and you can run it yourself")
txt(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.6),
    [[Rn("This isn't a mockup. The engine, the grading, the decisions — all built and working. A few numbers that prove it:", 15, FOG)]])
proof = [
    ("51 / 51", "tricky situations handled correctly"),
    ("100%", "always picks the best option in tests"),
    ("100%", "never breaks the safety & legal rules"),
    ("Live", "real website + working demo accounts"),
]
for i, (big, lab) in enumerate(proof):
    cx = Inches(0.7 + i * 3.02)
    rect(s, cx, Inches(2.5), Inches(2.82), Inches(1.5), PANEL, radius=True)
    rect(s, cx, Inches(2.5), Inches(2.82), Inches(0.1), ORANGE)
    txt(s, cx + Inches(0.25), Inches(2.72), Inches(2.4), Inches(0.6), [[Rn(big, 26, ORANGE, True)]])
    txt(s, cx + Inches(0.25), Inches(3.35), Inches(2.45), Inches(0.6), [[Rn(lab, 12, FOG)]], line_spacing=1.1)
screenshot_ph(s, Inches(0.7), Inches(4.35), Inches(11.9), Inches(2.0),
              "pnpm eval   &&   pnpm test:edge",
              "The full self-test report — everything green, reproducible from scratch on any laptop.")
footer(s, 16)

# ================================================================================
# 17 — why only Amazon / close
# ================================================================================
s = slide(NAVY)
accent_bar(s); kicker(s, "Why this is Amazon's to win")
title(s, "Amazon already owns the rails — we add the missing piece")
bullets(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.6), [
    ("The trucks, stations, and resale stores already exist. ", "ReLoop doesn't rebuild any of that — it just adds a brain at the very start."),
    ("The trust is already there. ", "Every item is Amazon-checked and guaranteed — not a stranger's promise. This is why it's not just another used-goods app."),
    ("It's a safe bet. ", "Worst case, it does exactly what happens today. Every good early decision on top of that is pure upside."),
], fs=15.5, gap=15)
rect(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.25), NAVY, line=ORANGE, line_w=1.5, radius=True)
txt(s, Inches(1.0), Inches(5.45), Inches(11.3), Inches(0.7),
    [[Rn("In one line:  ", 17, ORANGE, True),
      Rn("look at the item at your door, decide its best next home before it moves — and waste far less.", 17, WHITE, True)]],
    line_spacing=1.15)
footer(s, 17)

# ================================================================================
# 18 — final
# ================================================================================
s = slide(INK)
rect(s, 0, 0, EMU_W, Inches(0.14), ORANGE)
rect(s, 0, Inches(7.36), EMU_W, Inches(0.14), ORANGE)
txt(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(1.0), [[Rn("ReLoop", 52, WHITE, True)]])
txt(s, Inches(0.9), Inches(3.4), Inches(11.6), Inches(1.4),
    [[Rn("Grade at the doorstep.", 30, FOG, False)],
     [Rn("Decide before the item moves.", 30, ORANGE, True)]], line_spacing=1.2, space_after=4)
txt(s, Inches(0.9), Inches(5.5), Inches(11.6), Inches(0.6),
    [[Rn("“The landfill is a design flaw.”", 18, MUTED, False, FONT, True)]])
txt(s, Inches(0.9), Inches(6.3), Inches(11.6), Inches(0.5),
    [[Rn("reloop-woad.vercel.app", 15, CYAN, True)]])

prs.save("deck/ReLoop-Return-Pipeline.pptx")
print("saved deck/ReLoop-Return-Pipeline.pptx with", len(prs.slides._sldIdLst), "slides")
